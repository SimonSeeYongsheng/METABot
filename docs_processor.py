import PyPDF2
import logging
import json
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
import xlrd
from pptx import Presentation

logging.basicConfig(
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    level=logging.INFO
)

class Docs_processor:
    def __init__(self):
        pass

    def load_document(self, file_path: str, file_type: str):
        # Determine file type and process accordingly
        match file_type:
            case "pdf":
                try:
                    text = self.process_pdf(file_path)
                except Exception as e:
                    logging.error(f"Error loading PDF file: {e}")
                    raise ValueError(f"Failed to process PDF file: {file_path}")
            case "docx":
                try:
                    text = self.process_docx(file_path)
                except Exception as e:
                    logging.error(f"Error loading DOCX file: {e}")
                    raise ValueError(f"Failed to process DOCX file: {file_path}")
            case "json":
                try:
                    text = self.process_json(file_path)
                except Exception as e:
                    logging.error(f"Error loading JSON file: {e}")
                    raise ValueError(f"Failed to process JSON file: {file_path}")
            case "html" | "htm":
                try:
                    text = self.process_html(file_path)
                except Exception as e:
                    logging.error(f"Error loading HTML file: {e}")
                    raise ValueError(f"Failed to process HTML file: {file_path}")
            case "xml":
                try:
                    text = self.process_xml(file_path)
                except Exception as e:
                    logging.error(f"Error loading XML file: {e}")
                    raise ValueError(f"Failed to process XML file: {file_path}")
            case "xlsx":
                try:
                    text = self.process_xlsx(file_path)
                except Exception as e:
                    logging.error(f"Error loading XLSX file: {e}")
                    raise ValueError(f"Failed to process XLSX file: {file_path}")
            case "xls":
                try:
                    text = self.process_xls(file_path)
                except Exception as e:
                    logging.error(f"Error loading XLS file: {e}")
                    raise ValueError(f"Failed to process XLS file: {file_path}")
            case "ipynb":
                try:
                    text = self.process_ipynb(file_path)
                except Exception as e:
                    logging.error(f"Error loading IPYNB file: {e}")
                    raise ValueError(f"Failed to process IPYNB file: {file_path}")
            case "pptx":
                try:
                    text = self.process_pptx(file_path)
                except Exception as e:
                    logging.error(f"Error loading PPTX file: {e}")
                    raise ValueError(f"Failed to process PPTX file: {file_path}")
                
            case _:
                try:
                    text = self.process_text_file(file_path)
                except Exception as e:
                    logging.error(f"Error loading text file: {e}")
                    raise ValueError(f"Failed to process text file: {file_path}")
        return text

    def process_pdf(self, file_path):
        text = ""
        with open(file_path, "rb") as pdf_file:
            reader = PyPDF2.PdfReader(pdf_file)
            for page in reader.pages:
                text += page.extract_text()
        return text

    def process_docx(self, file_path):
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text

    def process_json(self, file_path):
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return json.dumps(data, indent=4)

    def process_html(self, file_path):
        with open(file_path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "html.parser")
        return soup.get_text(separator="\n")

    def process_xml(self, file_path):
        tree = ET.parse(file_path)
        root = tree.getroot()
        text = " ".join(root.itertext())
        return text

    def process_xlsx(self, file_path):
        wb = load_workbook(file_path, read_only=True)
        texts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        texts.append(str(cell))
        return "\n".join(texts)

    def process_xls(self, file_path):
        wb = xlrd.open_workbook(file_path)
        texts = []
        for sheet in wb.sheets():
            for row in range(sheet.nrows):
                row_values = sheet.row_values(row)
                texts.append(" ".join([str(cell) for cell in row_values if cell != ""]))
        return "\n".join(texts)

    def process_ipynb(self, file_path):
        with open(file_path, "r", encoding="utf-8") as f:
            notebook = json.load(f)
        texts = []
        for cell in notebook.get("cells", []):
            if cell.get("cell_type") in ["markdown", "code"]:
                cell_text = "".join(cell.get("source", []))
                texts.append(cell_text)
        return "\n".join(texts)

    def process_pptx(self, file_path):
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)

    def process_text_file(self, file_path):
        with open(file_path, "r", encoding="utf-8") as text_file:
            return text_file.read()